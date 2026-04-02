from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parent
REPORTS_DIR = PROJECT_ROOT / "artifacts" / "reports"
PLOTS_DIR = PROJECT_ROOT / "artifacts" / "plots"
SCREENSHOTS_DIR = PROJECT_ROOT / "artifacts" / "screenshots"
OUTPUT_PATH = REPORTS_DIR / "ids_final_delivery_presentation.pptx"

TITLE = "Multi-Class Intrusion Detection using Machine Learning: A Comparative and Explainable Approach"
SUBTITLE = "Hrushikesh Ramilla | Avadh Khandelwal | Himavarshith Reddy"


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = TITLE
    slide.placeholders[1].text = SUBTITLE


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    for index, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(24)


def add_image_slide(prs: Presentation, title: str, bullets: list[str], image_path: Path, left: float = 6.0) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.1), Inches(5.5))
    tf = tx_box.text_frame
    tf.clear()
    for index, bullet in enumerate(bullets):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(21)
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(left), Inches(1.4), width=Inches(6.8))


def add_two_image_slide(prs: Presentation, title: str, bullets: list[str], image_a: Path, image_b: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    tx_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.0), Inches(12.5), Inches(1.2))
    tf = tx_box.text_frame
    tf.clear()
    for index, bullet in enumerate(bullets):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(18)
    if image_a.exists():
        slide.shapes.add_picture(str(image_a), Inches(0.5), Inches(2.1), width=Inches(6.1))
    if image_b.exists():
        slide.shapes.add_picture(str(image_b), Inches(6.8), Inches(2.1), width=Inches(6.1))


def add_three_image_slide(prs: Presentation, title: str, bullets: list[str], images: list[Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    tx_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.9), Inches(12.2), Inches(1.0))
    tf = tx_box.text_frame
    tf.clear()
    for index, bullet in enumerate(bullets):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(17)
    slots = [(0.4, 2.0, 4.1), (4.6, 2.0, 4.1), (8.8, 2.0, 4.1)]
    for image_path, (left, top, width) in zip(images, slots):
        if image_path.exists():
            slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width))


def add_metric_table_slide(prs: Presentation, title: str, headers: list[str], rows: list[list[str]], notes: list[str], image_path: Path | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    rows_count = len(rows) + 1
    table = slide.shapes.add_table(rows_count, len(headers), Inches(0.4), Inches(1.1), Inches(6.0), Inches(4.2)).table
    for col_idx, header in enumerate(headers):
        table.cell(0, col_idx).text = header
    for row_idx, values in enumerate(rows, start=1):
        for col_idx, value in enumerate(values):
            table.cell(row_idx, col_idx).text = value
    notes_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(5.8), Inches(1.3))
    tf = notes_box.text_frame
    tf.clear()
    for index, note in enumerate(notes):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = note
        paragraph.level = 0
        paragraph.font.size = Pt(18)
    if image_path and image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(6.6), Inches(1.3), width=Inches(6.4))


def add_architecture_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Model Architecture"
    steps = [
        "Raw UNSW Features",
        "Cleaning + ID Removal",
        "Interaction Features",
        "Encoding + Selection",
        "SMOTE + Class Weights",
        "LGBM",
        "Calibration + Thresholds",
        "Prediction + SHAP",
    ]
    left = 0.5
    top = 2.0
    width = 1.45
    height = 1.0
    for index, step in enumerate(steps):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left + index * 1.55),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shape.text = step
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
        if index < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.CHEVRON,
                Inches(left + index * 1.55 + 1.35),
                Inches(top + 0.2),
                Inches(0.3),
                Inches(0.5),
            )
            arrow.text = ""
    notes = slide.shapes.add_textbox(Inches(0.7), Inches(4.5), Inches(11.8), Inches(1.6))
    tf = notes.text_frame
    tf.text = "Production path: cleaned inputs -> selected 75-feature space -> calibrated + thresholded LGBM -> multiclass prediction -> SHAP explanation"
    tf.paragraphs[0].font.size = Pt(22)


def build_deck() -> Path:
    model_metrics = load_json(REPORTS_DIR / "model_metrics.json")
    final_summary = load_json(REPORTS_DIR / "final_summary.json")
    data_profile = load_json(REPORTS_DIR / "data_profile.json")
    gui_validation = load_json(REPORTS_DIR / "gui_validation.json")

    final_metrics = final_summary["calibrated_final_model_metrics"]
    thresholds = final_summary.get("class_thresholds", {})
    phase_rows = [
        ["Phase 1 Baseline", "71.71%", "0.4341", "0.6982"],
        ["Phase 2 Production", "73.49%", "0.4992", "0.6890"],
        ["Phase 3 Production", f"{final_metrics['accuracy'] * 100:.2f}%", f"{final_metrics['macro_f1']:.4f}", f"{final_metrics['weighted_f1']:.4f}"],
    ]
    final_rows = [
        ["Accuracy", f"{final_metrics['accuracy'] * 100:.2f}%"],
        ["Macro F1", f"{final_metrics['macro_f1']:.4f}"],
        ["Weighted F1", f"{final_metrics['weighted_f1']:.4f}"],
        ["Final Model", final_summary["final_model"].upper()],
        ["Fuzzers F1", f"{final_metrics['classification_report']['fuzzers']['f1-score']:.4f}"],
        ["Fuzzers Recall", f"{final_metrics['classification_report']['fuzzers']['recall']:.4f}"],
    ]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_bullet_slide(
        prs,
        "Problem Statement",
        [
            "Intrusion Detection Systems classify malicious network traffic before damage spreads.",
            "UNSW-NB15 is highly imbalanced, multiclass, and noisy, which makes minority-class detection difficult.",
            "Goal: build a reproducible, explainable IDS pipeline and a usable GUI for demo and SaaS-style inference.",
        ],
    )
    add_image_slide(
        prs,
        "Dataset Overview",
        [
            f"Predefined training rows: {data_profile.get('train_rows', 0):,}",
            f"Predefined testing rows: {data_profile.get('test_rows', 0):,}",
            "After leakage-safe cleaning, duplicates were removed and the multiclass target remained attack_cat.",
            "Class distribution remains heavily skewed toward normal and exploits traffic.",
        ],
        PLOTS_DIR / "label_distribution.png",
    )
    add_image_slide(
        prs,
        "Initial Problems",
        [
            "ID leakage inflated apparent performance and hid true generalization gaps.",
            "Duplicate removal originally ran before dropping ID, so duplicates were effectively invisible.",
            "Feature-selection and imbalance handling were too weak for minority classes like fuzzers and worms.",
        ],
        PLOTS_DIR / "repo_benchmark_accuracy_comparison.png",
    )
    add_metric_table_slide(
        prs,
        "Baseline Performance",
        ["Stage", "Accuracy", "Macro F1", "Weighted F1"],
        [["Leakage-safe baseline", "71.71%", "0.4341", "0.6982"]],
        [
            "Phase 1 established the first reliable, leakage-free baseline after removing ID leakage and fixing duplicate handling.",
            "LGBM became the best baseline model under the unchanged ranking rule.",
        ],
        PLOTS_DIR / "model_metric_comparison.png",
    )
    add_metric_table_slide(
        prs,
        "Phase 1 Fixes",
        ["Fix", "Why It Mattered"],
        [
            ["Drop ID before features", "Removed leakage and reduced overfitting."],
            ["Correct duplicate order", "Exposed 26,387 train and 67,601 test duplicates."],
            ["Restore feature cap", "Re-enabled controlled selection for stable training."],
        ],
        [
            "These fixes rebuilt trust in the evaluation pipeline.",
            "They also changed the effective class distribution materially, especially for normal and generic traffic.",
        ],
        PLOTS_DIR / "missing_after_cleaning.png",
    )
    add_metric_table_slide(
        prs,
        "Phase 2 Improvements",
        ["Model", "Accuracy", "Macro F1", "Weighted F1"],
        [
            ["LGBM (offline)", "73.74%", "0.5248", "0.7010"],
            ["LGBM (production calibrated)", "73.49%", "0.4992", "0.6890"],
        ],
        [
            "Feature cap increased to 50, class weights were applied, SMOTE was added on training-only paths, and the final model was calibrated.",
            "Macro F1 improved materially, but fuzzers still underperformed in production.",
        ],
        PLOTS_DIR / "training_time_comparison.png",
    )
    add_metric_table_slide(
        prs,
        "Phase 3 Improvements",
        ["Model", "Accuracy", "Macro F1", "Weighted F1"],
        [
            ["LGBM (offline)", "75.34%", "0.5346", "0.7244"],
            ["LGBM (production thresholded)", "74.73%", "0.5269", "0.7125"],
        ],
        [
            "Added interaction features, raised feature cap to 75, reduced SMOTE aggressiveness, and tuned minority thresholds.",
            f"Tuned thresholds: {', '.join(f'{k}={v:.2f}' for k, v in thresholds.items())}.",
        ],
        PLOTS_DIR / "per_class_f1_heatmap.png",
    )
    add_metric_table_slide(
        prs,
        "Final Results",
        ["Metric", "Value"],
        final_rows,
        [
            "The current deployable model is a calibrated and thresholded LGBM.",
            "Performance improved materially across phases, but the 0.60+ Macro F1 target was not reached.",
        ],
        PLOTS_DIR / "confusion_matrix_lgbm.png",
    )
    add_architecture_slide(prs)
    add_three_image_slide(
        prs,
        "GUI Demo",
        [
            f"Dashboard validation passed: overview={gui_validation.get('overview_loaded')}, prediction={gui_validation.get('manual_prediction_rendered')}, CSV with id={gui_validation.get('csv_with_id_prediction_rendered')}.",
            "The GUI prefers the API when available and falls back to direct local inference when needed.",
            "Prediction flow supports manual single-record scoring, CSV batch upload, and SHAP-backed explanation views.",
        ],
        [
            SCREENSHOTS_DIR / "dashboard_overview.png",
            SCREENSHOTS_DIR / "dashboard_predict_form.png",
            SCREENSHOTS_DIR / "dashboard_prediction_result.png",
        ],
    )
    add_image_slide(
        prs,
        "Key Insights",
        [
            "Leakage removal mattered more than raw model complexity.",
            "Controlled SMOTE plus threshold tuning improved minority-class recovery without destabilizing the pipeline.",
            "Fuzzers remains the most difficult class because of overlap with benign behavior and limited discriminative signal.",
        ],
        PLOTS_DIR / "performance_evolution.png",
    )
    add_bullet_slide(
        prs,
        "Limitations",
        [
            "Production Macro F1 is 0.5269, so the 0.60+ target was not achieved.",
            "Fuzzers recall remains low even after threshold tuning and interaction features.",
            "The dataset still reflects strong class overlap and heavy skew toward normal traffic.",
        ],
    )
    add_bullet_slide(
        prs,
        "Future Work",
        [
            "Train a dedicated binary detector alongside the multiclass model for layered defense.",
            "Explore cost-sensitive boosting and richer temporal or flow-aggregation features.",
            "Add threshold search per deployment scenario and automated periodic model refresh.",
        ],
    )

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH


def main() -> None:
    output = build_deck()
    print(output)


if __name__ == "__main__":
    main()
